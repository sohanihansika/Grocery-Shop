from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(102, 126, 234)  # Purple-blue
SECONDARY_COLOR = RGBColor(118, 75, 162)  # Dark purple
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Title underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = PRIMARY_COLOR
    line.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.3

def add_code_slide(prs, title, code_text, description=""):
    """Add a slide with code block"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = PRIMARY_COLOR
    line.line.width = Pt(3)
    
    # Description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(0.6))
        desc_frame = desc_box.text_frame
        p = desc_frame.paragraphs[0]
        p.text = description
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
    
    # Code box
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2 if description else 1.5), Inches(8.4), Inches(4.8))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    p = code_frame.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.font.color.rgb = RGBColor(45, 45, 45)
    
    # Code background
    code_shape = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2 if description else 1.5), Inches(8.4), Inches(4.8))
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    code_shape.line.color.rgb = RGBColor(200, 200, 200)
    slide.shapes._spTree.remove(code_shape._element)
    slide.shapes._spTree.insert(2, code_shape._element)

# Slide 1: Title
add_title_slide(prs, "JavaScript Essentials", "Foundations of Modern JavaScript Engineering")

# Slide 2: What is JavaScript?
add_content_slide(prs, "What is JavaScript?", [
    "• High-level programming language - abstracts complex details",
    "• Single-threaded - executes one task at a time",
    "• Dynamically typed - data type decided at runtime",
    "• Interpreted language - executes code line by line",
    "• Automatic memory management - garbage collection built-in",
    "• Event-Driven - responds to user interactions in real-time"
])

# Slide 3: JavaScript Environments
add_content_slide(prs, "JavaScript Environments", [
    "CLIENT-SIDE:",
    "• Runs directly in the user's browser",
    "• Faster response time - no server communication needed",
    "• Brings interactivity to web pages",
    "",
    "SERVER-SIDE:",
    "• Executes on web servers",
    "• Accesses databases and file systems",
    "• Handles security features"
])

# Slide 4: Key Features
add_content_slide(prs, "Key Features of JavaScript", [
    "⚡ Client-Side Scripting - Fast response without server communication",
    "🎯 Versatile - From simple calculations to complex applications",
    "📡 Asynchronous - Fetches data without freezing the interface",
    "📚 Rich Ecosystem - React, Angular, Vue.js frameworks available"
])

# Slide 5: Limitations
add_content_slide(prs, "Limitations of JavaScript", [
    "🔓 Security Risks - Vulnerable to XSS attacks",
    "⚙️ Performance - Slower for complex tasks",
    "🧩 Complexity - Requires understanding of advanced concepts",
    "📋 Weak Type Checking - Dynamically typed without strict validation"
])

# Slide 6: Variables
add_content_slide(prs, "Variable Declaration (ES6+)", [
    "var - Function/Global Scoped (Old way)",
    "   • Can be redeclared and reassigned",
    "",
    "let - Block Scoped (Modern)",
    "   • Can be reassigned but NOT re-declared in same scope",
    "",
    "const - Block Scoped & Immutable",
    "   • Cannot be reassigned or re-declared"
])

# Slide 7: Code - Variables
add_code_slide(prs, "Variables Example", 
"let age = 25;\nconst PI = 3.14;\nvar name = 'John';\n\n// null vs undefined\nlet x = null;           // intentional\nlet y;                  // undefined",
"null: intentional absence | undefined: declared but not assigned")

# Slide 8: Arrow Functions
add_content_slide(prs, "ES6+ Features", [
    "Arrow Functions - Shorter syntax for function expressions",
    "   const add = (a, b) => a + b;",
    "",
    "Template Literals - Embed expressions and multi-line strings",
    "   console.log(`Hello ${name}, welcome!`);",
    "",
    "Default Parameters - Set default values for function parameters",
    "   const greet = (name = 'Guest') => { ... }"
])

# Slide 9: Destructuring
add_code_slide(prs, "Destructuring & Spread Operator",
"// Destructuring\nconst user = { name: 'Alice', age: 25 };\nconst { name, age } = user;\n\n// Spread Operator\nconst nums = [1, 2, 3];\nconst newNums = [...nums, 4, 5];",
"Easy extraction of values and combining arrays/objects")

# Slide 10: Modules
add_content_slide(prs, "Modules", [
    "✓ Organize Code - Break into focused, self-contained files",
    "✓ Prevent Naming Conflicts - Private scope by default",
    "✓ Improve Readability - Easier to navigate and understand",
    "✓ Reusability - Follow DRY (Don't Repeat Yourself) principle",
    "✓ Better Encapsulation - Variables private unless exported"
])

# Slide 11: Modules Code
add_code_slide(prs, "Modules Example",
"// math.js\nexport function add(a, b) { \n  return a + b; \n}\n\n// app.js\nimport { add } from './math.js';\nconsole.log(add(5, 3));  // 8",
"Export from one file, import to another")

# Slide 12: Async Programming
add_content_slide(prs, "Asynchronous Programming", [
    "Asynchronous operations complete concurrently without blocking main thread",
    "",
    "Common Use Cases:",
    "• API requests - fetching data from servers",
    "• File operations - reading/writing files",
    "• Event listeners - clicks, keyboard inputs",
    "• Database queries - retrieving data"
])

# Slide 13: Callbacks
add_content_slide(prs, "Callbacks", [
    "Functions passed as arguments, executed after a task completes",
    "",
    "✓ Widely used in async operations",
    "✓ Simple pattern for handling results",
    "",
    "✗ Callback Hell - Deeply nested callbacks become unreadable",
    "✗ Messy error handling with nested callbacks"
])

# Slide 14: Promises
add_content_slide(prs, "Promises Explained", [
    "Represents a value that may be available now, later, or never",
    "",
    "Three States:",
    "🔵 Pending - Task in initial state",
    "🟢 Fulfilled - Task completed successfully",
    "🔴 Rejected - Task failed with error"
])

# Slide 15: Promise Code
add_code_slide(prs, "Promises Example",
"const promise = new Promise((resolve, reject) => {\n  if (success) {\n    resolve(value);\n  } else {\n    reject(error);\n  }\n});\n\npromise\n  .then(result => console.log(result))\n  .catch(error => console.log(error));",
"Promise states and handling with .then() and .catch()")

# Slide 16: Promise Methods
add_content_slide(prs, "Promise Methods", [
    "Promise.all() - Wait for all, reject if any fails",
    "Promise.allSettled() - Wait for all, return all outcomes",
    "Promise.race() - Return result of first settled promise",
    "Promise.any() - Return first fulfilled promise",
    "Promise.resolve() - Resolve with given value",
    "Promise.reject() - Reject with given reason",
    "Promise.finally() - Cleanup code regardless of result"
])

# Slide 17: Async/Await
add_content_slide(prs, "Async/Await", [
    "Modern syntax making async code look synchronous",
    "",
    "async - Transforms function to return a Promise",
    "await - Pauses execution until Promise resolves",
    "",
    "Benefits:",
    "✓ More readable and maintainable code",
    "✓ Easier error handling with try/catch",
    "✓ Cleaner than .then() chains"
])

# Slide 18: Async/Await Code
add_code_slide(prs, "Async/Await Example",
"async function getData() {\n  try {\n    const response = await fetch('/api/data');\n    const data = await response.json();\n    console.log(data);\n  } catch (error) {\n    console.log('Error:', error);\n  }\n}",
"Async function with await and try/catch error handling")

# Slide 19: Scope
add_content_slide(prs, "Scope in JavaScript", [
    "Scope determines where variables can be accessed",
    "",
    "Global Scope - Accessible everywhere",
    "Function Scope - Accessible only within function",
    "Block Scope - let/const accessible only in { } block",
    "Lexical Scope - Inner functions access outer variables",
    "",
    "Note: var does not have block scope"
])

# Slide 20: Context (this)
add_content_slide(prs, "Context (this keyword)", [
    "Context refers to the value of 'this' in a function",
    "",
    "Global Context - Outside functions, this = global object",
    "Function Context - When called as method, this = object",
    "Constructor Context - With new keyword, this = new object",
    "Explicit Binding - Use call(), apply(), or bind()"
])

# Slide 21: Closures
add_content_slide(prs, "Closures", [
    "Inner function has access to outer function's variables",
    "Closures created every time a function is created",
    "",
    "Key Point:",
    "Inner function retains access to outer variables",
    "even after outer function finishes executing"
])

# Slide 22: Closures Code
add_code_slide(prs, "Closures Example",
"function outer() {\n  const outerVar = 'I am outer';\n  \n  function inner() {\n    console.log(outerVar);\n  }\n  \n  return inner;\n}\n\nconst closure = outer();\nclosure();  // 'I am outer'",
"Inner function retains access to outer variables")

# Slide 23: Closure Use Cases
add_content_slide(prs, "Closures - Use Cases", [
    "1. Data Privacy/Encapsulation",
    "   • Create private variables not accessible globally",
    "",
    "2. Function Factories",
    "   • Create functions with pre-configured behavior",
    "",
    "3. Partial Application",
    "   • Fix some arguments, create new functions"
])

# Slide 24: Closure Pitfalls
add_content_slide(prs, "Closures - Common Pitfalls", [
    "⚠️ Memory Leaks",
    "   • Excessive closures retain unnecessary references",
    "",
    "⚠️ Performance Overhead",
    "   • Overusing closures increases memory usage",
    "",
    "Best Practice: Use closures strategically"
])

# Slide 25: Conclusion
add_title_slide(prs, "Thank You!", "Master these JavaScript fundamentals to become an expert")

# Save presentation
prs.save('JavaScript_Essentials.pptx')
print("✅ PowerPoint presentation created successfully: JavaScript_Essentials.pptx")